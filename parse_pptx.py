from pptx import Presentation
import sys

def extract_text(pptx_path):
    prs = Presentation(pptx_path)
    text_runs = []
    for slide_number, slide in enumerate(prs.slides):
        text_runs.append(f"--- Slide {slide_number + 1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

if __name__ == "__main__":
    if len(sys.argv) > 1:
        text = extract_text(sys.argv[1])
        with open("pptx_content.txt", "w", encoding="utf-8") as f:
            f.write(text)
